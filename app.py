

import nltk
import os
import os
import nltk
import io
import streamlit as st
import pdfplumber
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from datetime import datetime
from PIL import Image
import torch
import torchvision.transforms as transforms
from torchvision import models
import re
from docx import Document
import groq

stopwords = {
    'i', 'me', 'my', 'myself', 'we', 'our', 'ours', 'ourselves', 'you', 'your', 'yours', 'yourself', 'yourselves',
    'he', 'him', 'his', 'himself', 'she', 'her', 'hers', 'herself', 'it', 'its', 'itself', 'they', 'them', 'their',
    'theirs', 'themselves', 'what', 'which', 'who', 'whom', 'this', 'that', 'these', 'those', 'am', 'is', 'are',
    'was', 'were', 'be', 'been', 'being', 'have', 'has', 'had', 'having', 'do', 'does', 'did', 'doing', 'a', 'an',
    'the', 'and', 'but', 'if', 'or', 'because', 'as', 'until', 'while', 'of', 'at', 'by', 'for', 'with', 'about',
    'against', 'between', 'into', 'through', 'during', 'before', 'after', 'above', 'below', 'to', 'from', 'up', 
    'down', 'in', 'out', 'on', 'off', 'over', 'under', 'again', 'further', 'then', 'once', 'here', 'there', 'when',
    'where', 'why', 'how', 'all', 'any', 'both', 'each', 'few', 'more', 'most', 'other', 'some', 'such', 'no', 'nor',
    'not', 'only', 'own', 'same', 'so', 'than', 'too', 'very', 's', 't', 'can', 'will', 'just', 'don', 'should',
    'now'
}
# Define paths and configurations
UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Load and configure the ResNet50 model for graph classification
graph_labels = ["bar", "line", "pie", "scatter"]
graph_model = models.resnet50(pretrained=True)
num_classes = len(graph_labels)
graph_model.fc = torch.nn.Linear(graph_model.fc.in_features, num_classes)
graph_model.eval()

# Define helper functions

    
def preprocess_text(text):
    text = re.sub(r'\s+', ' ', text)  # Replace multiple spaces/newlines with a single space
    text = text.strip()  # Remove leading and trailing whitespace
    words = text.split()
    filtered_words = [word for word in words if word.lower() not in stopwords]
    return ''.join(filtered_words)

def preprocess_image(image_path):
    input_image = Image.open(image_path).convert("RGB")
    preprocess = transforms.Compose([
        transforms.Resize(256),
        transforms.CenterCrop(224),
        transforms.ToTensor(),
        transforms.Normalize(mean=[0.485, 0.456, 0.406], std=[0.229, 0.224, 0.225]),
    ])
    input_tensor = preprocess(input_image)
    return input_tensor.unsqueeze(0)

def extract_images_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    image_files = []
    for slide_num, slide in enumerate(prs.slides):
        for shape_num, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = io.BytesIO(image.blob)
                image_file_name = f"slide_{slide_num+1}_image_{shape_num+1}.{image.ext}"
                with open(image_file_name, "wb") as f:
                    f.write(image_bytes.read())
                image_files.append(image_file_name)
    return image_files

def extract_text_from_pdf(pdf_path):
    with pdfplumber.open(pdf_path) as pdf:
        extracted_text = [preprocess_text(page.extract_text() or "") for page in pdf.pages]
        return "\n".join(extracted_text)
        #return "\n".join(page.extract_text() for page in pdf.pages)

def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    return " ".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

def read_word_file(file_path):
    doc = Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)

def summ(text, chunk_size=1000):
    summaries = []
    api_key = "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"
    client = groq.Client(api_key=api_key)

    for chunk in chunk_text(text, 1000):
        response = client.chat.completions.create(
            messages=[
                {
                    "role": "user",
                    "content": "Please summarize the following text:\n\n" + chunk,
                }
            ],
            model="llama3-8b-8192",
        )
        summaries.append(response.choices[0].message.content)

    combined_summary = " ".join(summaries)
    return combined_summary

#def answer_question(question, context):
#   api_key = "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"
#    client = groq.Client(api_key=api_key)
#   response = client.chat.completions.create(
#       messages=[
#           {
#               "role": "user",
#               "content": context + "\n\nQuestion: " + question,
#           }
#       ],
#       model="llama3-8b-8192",
#   )
#   return response.choices[0].message.content

def chunk_text(text, chunk_size=1000):
    """Splits the text into smaller chunks."""
    for i in range(0, len(text), chunk_size):
        yield text[i:i + chunk_size]

def answer_question(question, context):
    api_key = "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"
    client = groq.Client(api_key=api_key)
    best_answer = ""
    best_score = float('-inf')

    for chunk in chunk_text(context, chunk_size=1000):
        try:
            response = client.chat.completions.create(
                messages=[
                    {
                        "role": "user",
                        "content": chunk + "\n\nQuestion: " + question,
                    }
                ],
                model="llama3-8b-8192",
            )
            answer = response.choices[0].message.content
            
            # Simple scoring mechanism, customize as needed
            score = len(answer)  # This could be based on answer length, specific keywords, etc.
            if score > best_score:
                best_score = score
                best_answer = answer
        except Exception as e:
            st.error(f"Error: {e}")

    return best_answer
# Streamlit UI
st.title("Chatbot - ChatG")
st.write("Upload a PDF, PPT, or Word document, and ask questions about its content.")

uploaded_file = st.file_uploader("Choose a file", type=["pdf", "ppt", "pptx", "doc", "docx"])

if uploaded_file:
    # Save the uploaded file
    file_path = os.path.join(UPLOAD_FOLDER, uploaded_file.name)
    with open(file_path, "wb") as f:
        f.write(uploaded_file.getbuffer())

    if uploaded_file.name.endswith(".pdf"):
        text = extract_text_from_pdf(file_path)
    elif uploaded_file.name.endswith(".ppt") or uploaded_file.name.endswith(".pptx"):
        text = extract_text_from_ppt(file_path)
    elif uploaded_file.name.endswith(".doc") or uploaded_file.name.endswith(".docx"):
        text = read_word_file(file_path)
    
    preprocessed_text = preprocess_text(text)
    st.write("Text extracted from the document:")
    st.text(preprocessed_text)
    
    summary_option = st.checkbox("Summarize the document?")
    if summary_option:
        with st.spinner("Summarizing..."):
            summary = summ(preprocessed_text)
        st.write("Summary:")
        st.markdown(
            f"""
            <div style="height:200px;overflow-y:scroll">
                {summary}
            </div>
            """,
            unsafe_allow_html=True
        )
    
    question = st.text_input("Ask a question about the document:")
    if question:
        with st.spinner("Answering..."):
            answer = answer_question(question, preprocessed_text)
        st.write("Answer:")
        st.markdown(
            f"""
            <div style="height:200px;overflow-y:scroll">
                {answer}
            </div>
            """,
            unsafe_allow_html=True
        )
    
    # question = st.text_input("Ask a question about the document:")
    # if question:
    #     with st.spinner("Answering..."):
    #         answer = answer_question(question, preprocessed_text)
    #     st.write("Answer:")
    #     st.text(answer)
